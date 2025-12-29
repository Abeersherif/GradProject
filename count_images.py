from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

file_path = r"c:\Users\Administrator\Desktop\Agents\GOLD-2026_teaching-slide-set-FINAL-v1.1-11Nov2026_Read-Only.pptx"
prs = Presentation(file_path)

image_count = 0
total_shapes = 0

def match_shape(shape):
    global image_count, total_shapes
    total_shapes += 1
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        image_count += 1
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            match_shape(s)

for slide in prs.slides:
    for shape in slide.shapes:
        match_shape(shape)

print(f"Total shapes: {total_shapes}")
print(f"Total images: {image_count}")
