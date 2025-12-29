from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

file_path = r"c:\Users\Administrator\Desktop\Agents\GOLD-2026_teaching-slide-set-FINAL-v1.1-11Nov2026_Read-Only.pptx"
output_file = r"c:\Users\Administrator\Desktop\Agents\ppt_content.txt"

if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    exit(1)

prs = Presentation(file_path)

def get_shape_text(shape):
    text = ""
    if hasattr(shape, "text") and shape.text.strip():
        text += shape.text.strip() + "\n"
    
    if shape.has_table:
        for row in shape.table.rows:
            row_text = []
            for cell in row.cells:
                if cell.text_frame and cell.text_frame.text.strip():
                    row_text.append(cell.text_frame.text.strip())
            if row_text:
                text += " | ".join(row_text) + "\n"
                
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            text += get_shape_text(s)
            
    return text

with open(output_file, "w", encoding="utf-8") as f:
    for i, slide in enumerate(prs.slides):
        f.write(f"\n--- Slide {i+1} ---\n")
        
        # Check for images to note their presence
        image_count = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
        if image_count > 0:
            f.write(f"[Contains {image_count} images]\n")

        # Extract text
        for shape in slide.shapes:
            text = get_shape_text(shape)
            if text.strip():
                f.write(text + "\n")

print(f"Dumped content to {output_file}")
