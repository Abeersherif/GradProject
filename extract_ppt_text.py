from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

file_path = r"c:\Users\Administrator\Desktop\Agents\GOLD-2026_teaching-slide-set-FINAL-v1.1-11Nov2026_Read-Only.pptx"

if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    exit(1)

prs = Presentation(file_path)

def iter_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape

def get_text_from_shape(shape):
    text = ""
    if hasattr(shape, "text"):
        text += shape.text + "\n"
    if shape.has_table:
        for row in shape.table.rows:
            row_text = []
            for cell in row.cells:
                if cell.text_frame:
                    row_text.append(cell.text_frame.text)
            text += " | ".join(row_text) + "\n"
    return text

for i, slide in enumerate(prs.slides):
    print(f"--- Slide {i+1} ---")
    
    # Slide Content
    for shape in iter_shapes(slide.shapes):
        text = get_text_from_shape(shape)
        if text.strip():
            print(f"[Content]: {text.strip()}")

    # Speaker Notes
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        notes = slide.notes_slide.notes_text_frame.text
        if notes.strip():
            print(f"[Notes]: {notes.strip()}")
            
    print("\n")
